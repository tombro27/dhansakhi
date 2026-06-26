#!/usr/bin/env python3
"""Fill the provided IDBI Innovate PPT template with DhanSakhi content, then it is
exported to PDF with LibreOffice. Content sits BELOW each slide's title bar so it
never overlaps the template's header/title."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

ROOT = "/home/imart/indiamart/hackathons/IDBI_Innovate_26"
TEMPLATE = os.path.join(ROOT, "deck", "Prototype Submission Deck _ IDBI Innovate.pptx")
OUT = os.path.join(ROOT, "deck", "DhanSakhi_IDBI_Innovate_Deck.pptx")
SHOTS = os.path.join(ROOT, "deck", "screenshots")

NAVY = RGBColor(0x0A, 0x25, 0x40)
INK = RGBColor(0x2B, 0x3A, 0x4D)
TEAL = RGBColor(0x00, 0x6D, 0x54)

prs = Presentation(TEMPLATE)
s = prs.slides


def body(slide, lines, top=1.55, left=0.5, width=9.0, height=3.8, size=12.5, color=INK):
    """Add a content text box below the title bar."""
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        bullet = line.startswith("• ")
        r = p.add_run()
        r.text = ("•  " + line[2:]) if bullet else line
        r.font.size = Pt(size)
        r.font.color.rgb = TEAL if (line and not bullet and line.endswith(":")) else color
        if line and not bullet and (line.endswith(":") or line.endswith("?")):
            r.font.bold = True
        p.space_after = Pt(5)
    return tb


def add_img(slide, path, left, top, width):
    if os.path.exists(path):
        slide.shapes.add_picture(path, Inches(left), Inches(top), width=Inches(width))


# ---- Slide 1: fill the existing Team-Details box (below the banner) ----
tb = next(sh for sh in s[0].shapes if sh.has_text_frame and "Team" in sh.text_frame.text)
tf = tb.text_frame
tf.clear()
rows = [
    ("Team Details", 18, True, NAVY),
    ("", 8, False, INK),
    ("Team name:  Team DhanSakhi", 14, False, INK),
    ("Team leader:  Prakhar Singh Tomar", 14, False, INK),
    ("Problem Statement:  Digital Wealth Management", 14, False, INK),
    ("Solution:  DhanSakhi — an AI Robo-Advisor for IDBI Bank", 14, False, TEAL),
]
for i, (txt, sz, bold, col) in enumerate(rows):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    r = p.add_run(); r.text = txt
    r.font.size = Pt(sz); r.font.bold = bold; r.font.color.rgb = col

# ---- 2 Brief about the idea ----
body(s[1], [
    "DhanSakhi gives every IDBI customer a personal robo-advisor:",
    "• One-tap RBI Account Aggregator consent pulls ALL holdings — bank, MF, demat, EPF, NPS, insurance.",
    "• AI builds goal-based plans and an explainable, suitability-aware portfolio.",
    "• Every recommendation carries a plain-language 'why' — no black boxes.",
    "• A vernacular conversational advisor (Google Gemini) answers anything, in 6 languages.",
    "• Goals map to real IDBI + LIC products (FD, MF, NPS via IDBI POP, LIC insurance/annuity).",
    "",
    "Why now: AA covers 2.61B accounts / 252.9M users (Dec 2025); IDBI already owns the distribution",
    "shelf and LIC (~49% owner). Quality advice is HNI-only today — DhanSakhi democratises it.",
], top=1.55)

# ---- 3 Opportunities ----
body(s[2], [
    "How is it different?",
    "• AA-native aggregation — no CAS/PDF imports that INDmoney & HDFC SmartWealth still need.",
    "• Explainable AI rationale vs black-box curation or rule-based recommendations.",
    "• Vernacular advisor in 6 languages — depth almost no competitor offers.",
    "How will it solve the problem?",
    "• One 360° view; AI risk-profile → target allocation → rebalancing with reasons; goals funded",
    "  by the right IDBI/LIC product.",
    "USP:  'Trust of a bank + intelligence of a super-app' — the intersection no standalone app or",
    "bank tool owns.",
], top=2.05, size=12)

# ---- 4 Features ----
body(s[3], [
    "• KYC-lite onboarding (DigiLocker/Aadhaar-style, paperless)",
    "• RBI Account Aggregator consent artefact (purpose, scope, expiry, one-tap revoke)",
    "• 360° net-worth dashboard (bank + FD + MF + demat + EPF + NPS + LIC)",
    "• AI risk profiling → Conservative / Moderate / Aggressive, explained",
    "• Goal-based planning (retirement, education, home, emergency) with SIP-to-goal math",
    "• Explainable portfolio advice + rebalancing ('why this, why now')",
    "• Vernacular conversational advisor (Gemini, 6 languages, suitability-aware)",
    "• IDBI + LIC cross-sell engine (next-best-product mapped to goal gaps)",
    "• Trust & compliance layer (SEBI guardrails, DPDP consent, no-PII)",
    "• Cross-sell ROI model (benchmark-anchored revenue projection for IDBI)",
], top=1.55, size=12.5)

# ---- 5 Process flow ----
body(s[4], [
    "Golden path (end-to-end):",
    "1) KYC-lite onboarding  →  2) AA consent & 360° aggregation  →  3) AI risk profile  →",
    "4) Goal-based plan  →  5) Explainable advice & rebalance  →  6) Vernacular advisor  →  7) IDBI/LIC cross-sell",
    "",
    "Actors:  Customer · DhanSakhi AI · RBI Account Aggregator (FIP→FIU) · IDBI core & product catalog · LIC · SEBI-RIA (human-in-loop).",
    "Trigger:  customer opens DhanSakhi inside IDBI net-banking and grants one consent to see and plan their whole financial life.",
    "Outcome:  a funded, explainable plan + the single best next product + an advisor they can ask anything, in their language.",
], top=1.55)

# ---- 6 Wireframes / mock (screenshots) ----
body(s[5], ["Hero flow — RBI Account Aggregator consent artefact (left) and AI risk-profiling (right):"], top=1.5, size=12, color=NAVY)
add_img(s[5], os.path.join(SHOTS, "02-onboarding-consent.png"), 0.5, 1.95, 4.4)
add_img(s[5], os.path.join(SHOTS, "09-onboarding-risk.png"), 5.1, 1.95, 4.4)

# ---- 7 Architecture ----
body(s[6], [
    "Customer (web / mobile, inside IDBI net-banking)",
    "        ▼ HTTPS",
    "Next.js 16 App Router on Vercel — Onboarding · Dashboard · Goals · Explainable Advice · Advisor · Cross-sell",
    "(React 19 + Tailwind + Recharts)",
    "        ▼ server route handlers",
    "AA Connector (sim. FIP→FIU consent)  |  Recommendation Engine (deterministic TS)  |",
    "Gemini API (advisor + explainability, no-key fallback)  |  IDBI + LIC product catalog",
    "        ▼",
    "Synthetic data store (no real PII).   POC path → IDBI core banking + live AA FIU + SEBI-RIA review.",
], top=1.55, size=12.5)

# ---- 8 Technologies ----
body(s[7], [
    "Frontend:  Next.js 16 (App Router), React 19, TypeScript, Tailwind CSS v4, Recharts",
    "AI & logic:  Google Gemini (2.5 Flash) via Google Gen AI SDK; deterministic TS engine",
    "  (risk, allocation, goals, rebalancing, cross-sell); natural-language explainability layer",
    "Integrations & infra:  RBI Account Aggregator (FIU), IDBI core + LIC/NPS catalog,",
    "  DigiLocker/Aadhaar KYC, Vercel (deploy + edge)",
    "Compliance:  SEBI IA guardrails, DPDP Act 2023, synthetic/no-PII, Bhashini-ready vernacular",
], top=1.6)

# ---- 9 Implementation cost ----
body(s[8], [
    "Prototype (done):  this app — ≈ ₹0 (open-source + free tiers, incl. Gemini free tier)",
    "Pilot / Sandbox (3–4 months):  live AA FIU, IDBI core + LIC APIs, RIA workflow, security review — ₹35–55 L",
    "Production (6–9 months):  scale infra, vernacular (Bhashini), DPO/DPIA, mobile app — ₹1.2–1.8 Cr",
    "Run-rate:  cloud + Gemini API + AA consent fees + support — ₹20–40 / active user / year",
    "",
    "Mostly variable, usage-based costs; AA + Gemini (free tier) + Vercel keep fixed costs low.",
    "Payback driven by cross-sell uplift (see slide 11).",
], top=1.6)

# ---- 10 Snapshots (screenshots grid) ----
body(s[9], ["A working, deployed product (idbi-innovate-26.vercel.app) — not mockups:"], top=1.45, size=12, color=NAVY)
grid = ["03-dashboard.png", "04-advice.png", "05-advisor.png", "06-goals.png", "07-cross-sell.png", "08-compliance.png"]
for idx, name in enumerate(grid):
    col, row = idx % 3, idx // 3
    add_img(s[9], os.path.join(SHOTS, name), 0.5 + col * 3.1, 1.8 + row * 1.82, 2.9)

# ---- 11 Performance / benchmarking ----
body(s[10], [
    "Prototype metrics:",
    "• 360° aggregation (6–8 institutions, simulated):  < 2 s",
    "• Explainability coverage of recommendations:  100%",
    "• Advisor reply latency:  ~1–3 s (Gemini) / < 50 ms (fallback)",
    "• Languages supported: 6   •   Uptime without API key (graceful fallback): 100%",
    "Projected impact (industry benchmarks — labelled as projections, not traction):",
    "• Cross-sell conversion lift ~6× (2% → ~12%)   •   Revenue uplift 25–40%",
    "• Illustrative incremental revenue ₹4.8 Cr / campaign   •   Payback < 18 months",
    "• Vernacular PoC (literature): +41% task completion, +86% session length vs English-only",
], top=1.55, size=12)

# ---- 12 Future development ----
body(s[11], [
    "Near term:  live AA FIU integration; IDBI core + LIC/NPS APIs; SEBI-RIA review workflow;",
    "  mobile app + WhatsApp advisor.",
    "Mid term:  Bhashini voice in 22 languages; auto-rebalance & SIP execution (EOP);",
    "  tax-harvesting & family dashboard; proactive nudges.",
    "Long term:  predictive churn & life-event triggers; RM co-pilot (B2B2C); open-finance",
    "  marketplace; financial inclusion at PMJDY scale.",
    "",
    "Winners may access IDBI's sandbox & POC environment — DhanSakhi drops onto IDBI's existing rails with low integration lift.",
], top=1.6)

# ---- 13 Links ----
body(s[12], [
    "GitHub Public Repository:",
    "   https://github.com/tombro27/dhansakhi",
    "Final Product Link (live deployment):",
    "   https://idbi-innovate-26.vercel.app",
    "Demo Video Link (3 minutes):",
    "   https://drive.google.com/file/d/1-ZobKRuBWrzYHAqD9f_Xko-KFBlIHDsC/view",
], top=2.45, size=14)

# ---- 14 Closing ----
if len(s) > 13:
    body(s[13], [
        "Thank you.",
        "",
        "DhanSakhi — the trust of a bank, with the intelligence of a super-app.",
        "Explainable, vernacular, Account-Aggregator-native wealth advice — ready for IDBI.",
    ], top=2.0, size=16, color=NAVY)

prs.save(OUT)
print("Saved", OUT, "with", len(s), "slides")
